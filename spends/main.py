from model.month import Month
from model.year import Year
from model.model import Model
from graphs.graphs import *
import json
from pptx import Presentation


# Formaters:
def format_numbers(list):
    result = []
    for x in list:
        result.append(format_number(x))
    return result

def format_number(number):
    return round(number,2)


try:
    model = Model()
    model.load_spends('datos.csv')

    # Obtencion datos:
    # Total
    canidad_fide = format_number(model.get_total_by_type_owner('total', 'fide'))
    canidad_comun = format_number(model.get_total_by_type_owner('total', 'comun'))
    canidad_total = canidad_comun + canidad_fide
    # Evolucion año actual
    total = model.get_monthly_ac_by_type_year('total', 2017)
    comun = model.get_monthly_ac_by_type_owner_year('total', 'comun', 2017)
    fide = model.get_monthly_ac_by_type_owner_year('total', 'fide', 2017)
    genarate_yearly_comparative(format_numbers(total), format_numbers(comun),format_numbers(fide), 2017)
    # Presentacion
    prs = Presentation()
    titulo = prs.slide_layouts[0]
    _titulo = prs.slides.add_slide(titulo)
    _titulo_title = _titulo.shapes.title
    _titulo_subtitle = _titulo.placeholders[1]

    _titulo_title.text = "Análisis Cuentas"
    _titulo_subtitle.text = "- " + str(2017) + " -"

    # Resumen anual
    resumen_anual = prs.slide_layouts[1]
    _resumen_anual = prs.slides.add_slide(resumen_anual)
    _resumen_anual_shapes = _resumen_anual.shapes
    _resumen_anual_title = _resumen_anual_shapes.title
    _resumen_anual_body = _resumen_anual_shapes.placeholders[1]
    _resumen_anual_title.text = 'Resumen cuentas ' + str(2017)
    tf = _resumen_anual_body.text_frame
    tf.text = 'Total cuenta: ' + str(canidad_total)  + ' €'
    p = tf.add_paragraph()
    p.level = 1
    p = tf.add_paragraph()
    p.text = 'Cuenta comun: ' + str(canidad_comun)  + ' €'
    p.level = 1
    p = tf.add_paragraph()
    p.level = 1
    p = tf.add_paragraph()
    p.text = 'Cuenta fide: ' + str(canidad_fide)  + ' €'
    p.level = 1

    # Se guarda el ppt
    prs.save('resultados.pptx')
    #evolucion entre años
    #Tabla gastos comparando mes actual vs media meses %
    #print((model.get_monthly_by_type_year('total', 2017)))
    #Tabla gastos medios vs actual %

    #Fide y comun:
    #Total
    #Evolucion año actual
    #print(model.get_monthly_ac_by_type_owner_year('total', 'comun', 2017))
    #print((model.get_monthly_by_type_owner_year('total', 'fide', 2017)))
    #evolucion entre años


except Exception as ex:
    print("-- error --")  
    print(ex)  